"""Generate the Eco-Matic presentation with a modern layered design and animated objects."""
from __future__ import annotations

from pathlib import Path
from typing import Any, Iterable, List, Sequence, cast

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "EcoMatic_Presentation.pptx"

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
CONTENT_LEFT = Inches(1.0)
CONTENT_WIDTH = SLIDE_WIDTH - Inches(2.0)

SHAPE = cast(Any, MSO_AUTO_SHAPE_TYPE)

PRIMARY = RGBColor(14, 25, 38)
ACCENT = RGBColor(74, 199, 178)
ACCENT_ALT = RGBColor(118, 141, 255)
SECONDARY = RGBColor(255, 196, 120)
TEXT = RGBColor(34, 40, 48)
LIGHT = RGBColor(249, 251, 254)
CANVAS_TINT = RGBColor(225, 230, 237)
MUTED = RGBColor(126, 137, 153)

SUBTITLE_TAG = "C# Console Application Project Documentation"
NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

SECTION_CONTENT = [
    {
        "number": 1,
        "title": "Project Overview",
        "subtitle": "Purpose • Features • Audience • Stack",
        "bullets": [
            "Purpose: Console-based kiosk that rewards recycling to spotlight SDG 12.",
            "Features: Dual customer/admin flows, CSV persistence, Spectre.Console dashboards, UTC logging.",
            "Target Audience: Instructors, classmates, and sustainability advocates evaluating the concept.",
            "Technology Stack: C# on .NET 9 with Spectre.Console, System.IO, and ANSI-styled output.",
        ],
    },
    {
        "number": 2,
        "title": "Requirements",
        "subtitle": "Software • Installation • Hardware",
        "bullets": [
            "Software: .NET 9 SDK, Spectre.Console package, and VS Code/VS 2022.",
            "Installation: `dotnet restore && dotnet build eco-matic/eco-matic.csproj` to prime data.",
            "Run: `dotnet run --project eco-matic/eco-matic.csproj` from the repository root.",
            "System: Any Windows laptop capable of running console applications.",
        ],
    },
    {
        "number": 3,
        "title": "File Handling Overview",
        "subtitle": "Inventory • Event Log • Safeguards",
        "bullets": [
            "Inventory CSV tracks slot type, name, price, stock, calories/volume.",
            "Event log CSV chronicles purchases, recycle credits, and admin actions with timestamps.",
            "Operations: load arrays on startup, rewrite inventory after edits, append log entries instantly.",
            "Safeguards: Header validation, numeric checks, and auto-regeneration for missing files.",
        ],
    },
    {
        "number": 4,
        "title": "Code Structure",
        "subtitle": "Classes • Methods • Modularity",
        "bullets": [
            "`EcoMatic` coordinates menus, balance handling, inventory updates, and logging.",
            "`VendingItem` base + Snack/Drink/Misc overrides deliver flavor text and CSV serialization.",
            "Helpers (`Write`, `TransactionTracker`, `RecycleTracker`, `SalesReport`) isolate responsibilities.",
            "Modularity: Constants define slot/stock limits to keep the single-file layout readable.",
        ],
    },
    {
        "number": 5,
        "title": "User Interface",
        "subtitle": "Design • Input/Output • Errors",
        "bullets": [
            "Spectre.Console tables and color-coded stock dots create a dashboard look.",
            "Customer path: insert bills → examine/buy items → recycle → receive change + receipt.",
            "Admin path: password gate for restock, add/remove items, view/clear logs, run sales report.",
            "Error guidance: the `Write` helper highlights invalid bills, IDs, and ranges.",
        ],
    },
    {
        "number": 6,
        "title": "Challenges & Solutions",
        "subtitle": "Development Hurdles",
        "bullets": [
            "CSV integrity: validation + auto-rebuild routines prevent corruption during demos.",
            "Array constraints: trackers reset every session to avoid overflow and stale data.",
            "Single-file rule: helper classes and regions keep logic modular despite the constraint.",
            "Storytelling: UX copy keeps the SDG recycling narrative front-and-center.",
        ],
    },
    {
        "number": 7,
        "title": "Testing",
        "subtitle": "Cases • Results • Limitations",
        "bullets": [
            "Cases: normal purchases, insufficient funds, sold-out items, recycle bounds, admin tasks.",
            "Evidence: reviewed CSV diffs and event log lines after each scenario.",
            "Tooling: `dotnet build` before demos guarantees a clean compile state.",
            "Limitations: testing is still manual; scripted regressions are planned.",
        ],
    },
    {
        "number": 8,
        "title": "Future Enhancements",
        "subtitle": "Planned Features • Performance",
        "bullets": [
            "Refactor by splitting `Program.cs` and adopting collections for scalability.",
            "Analytics: configurable recycle catalog plus multi-day dashboards.",
            "UX: GUI or web client for kiosk-ready deployments and remote admin control.",
            "Performance: inventory valuation, restock forecasting, and caching strategies.",
        ],
    },
    {
        "number": 9,
        "title": "Conclusion",
        "subtitle": "Reflection • Takeaways",
        "bullets": [
            "Learned to blend OOP, file I/O, and UX storytelling inside a console context.",
            "Validated that sustainability narratives can live within technical deliverables.",
            "Project already ships with MIT license, docs, and a ready-to-demo build.",
            "Next steps: modular refactor plus analytics-driven roadmap.",
        ],
    },
]

APPENDIX_CONTENT = {
    "number": 10,
    "title": "Appendix",
    "subtitle": "Source Code & References",
    "left_header": "Source Code",
    "left_items": [
        "Repo: Eco-Matic (branch feature/eco-matic-refactoring).",
        "Key file: `eco-matic/Program.cs` (single-file implementation).",
        "Data: `eco-matic/data/inventory.csv` + `eventLog.csv`.",
    ],
    "right_header": "References",
    "right_items": [
        "Spectre.Console docs for tables, markup, and styling.",
        ".NET documentation for file handling and console APIs.",
        "UN SDG 12 resources inspiring the recycle-for-credit concept.",
    ],
}


def apply_transition(slide) -> None:
    slide_element = slide._element
    transition = slide_element.find(qn("p:transition"))
    if transition is None:
        transition = OxmlElement("p:transition")
        c_sld = slide_element.find(qn("p:cSld"))
        if c_sld is None:
            slide_element.append(transition)
        else:
            slide_element.insert(list(slide_element).index(c_sld) + 1, transition)
    else:
        for child in list(transition):
            transition.remove(child)
    transition.set("spd", "slow")
    transition.append(OxmlElement("p:push"))


def get_shape_id(shape) -> str | None:
    ids = shape._element.xpath(".//p:cNvPr/@id")
    return ids[0] if ids else None


def apply_object_animations(slide, shapes: Iterable[Any], delay_increment: int = 150) -> None:
    filtered = [s for s in shapes if s is not None]
    if not filtered:
        return

    slide_element = slide._element
    existing = slide_element.find(qn("p:timing"))
    if existing is not None:
        slide_element.remove(existing)

    timing = OxmlElement("p:timing")
    tn_lst = OxmlElement("p:tnLst")
    timing.append(tn_lst)
    par = OxmlElement("p:par")
    tn_lst.append(par)

    root = OxmlElement("p:cTn")
    root.set("id", "1")
    root.set("dur", "indefinite")
    root.set("restart", "never")
    root.set("nodeType", "tmRoot")
    par.append(root)

    child_root = OxmlElement("p:childTnLst")
    root.append(child_root)

    seq = OxmlElement("p:seq")
    seq.set("concurrent", "1")
    seq.set("nextAc", "seek")
    child_root.append(seq)

    seq_ctn = OxmlElement("p:cTn")
    seq_ctn.set("id", "2")
    seq_ctn.set("dur", "indefinite")
    seq.append(seq_ctn)

    seq_child = OxmlElement("p:childTnLst")
    seq_ctn.append(seq_child)

    anim_id = 3
    delay = 0
    for shape in filtered:
        spid = get_shape_id(shape)
        if not spid:
            continue

        par_node = OxmlElement("p:par")
        seq_child.append(par_node)

        ctn = OxmlElement("p:cTn")
        ctn.set("id", str(anim_id))
        anim_id += 1
        ctn.set("dur", "500")
        ctn.set("fill", "hold")
        st_cond_lst = OxmlElement("p:stCondLst")
        cond = OxmlElement("p:cond")
        cond.set("delay", str(delay))
        cond.set("evt", "onBegin")
        st_cond_lst.append(cond)
        ctn.append(st_cond_lst)

        child_list = OxmlElement("p:childTnLst")
        ctn.append(child_list)
        anim = OxmlElement("p:animEffect")
        anim.set("transition", "in")
        anim.set("filter", "fade")
        anim.set("advClick", "0")
        anim.set("lvl", "0")
        child_list.append(anim)

        c_bhvr = OxmlElement("p:cBhvr")
        anim.append(c_bhvr)
        c_bhvr_c_tn = OxmlElement("p:cTn")
        c_bhvr_c_tn.set("id", str(anim_id))
        anim_id += 1
        c_bhvr_c_tn.set("dur", "500")
        c_bhvr_c_tn.set("fill", "hold")
        c_bhvr.append(c_bhvr_c_tn)

        tgt_el = OxmlElement("p:tgtEl")
        sp_tgt = OxmlElement("p:spTgt")
        sp_tgt.set("spid", spid)
        tgt_el.append(sp_tgt)
        c_bhvr.append(tgt_el)

        par_node.append(ctn)
        delay += delay_increment

    slide_element.append(timing)


def style_background(slide, accent: str = "right") -> List[Any]:
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = 90
    stops = fill.gradient_stops
    stops[0].color.rgb = LIGHT
    stops[0].position = 0.0
    stops[-1].color.rgb = CANVAS_TINT
    stops[-1].position = 1.0

    # Minimal accent line for subtle structure
    shapes: List[Any] = []
    line_x = SLIDE_WIDTH - Inches(0.6) if accent == "right" else Inches(0.5)
    accent_line = slide.shapes.add_shape(
        SHAPE.RECTANGLE,
        line_x,
        Inches(0.2),
        Inches(0.1),
        SLIDE_HEIGHT - Inches(0.4),
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT if accent == "right" else ACCENT_ALT
    accent_line.line.fill.background()
    accent_line.fill.transparency = 0.25
    shapes.append(accent_line)

    return shapes


def add_header(slide, text: str = SUBTITLE_TAG):
    header = slide.shapes.add_textbox(CONTENT_LEFT, Inches(0.25), CONTENT_WIDTH, Inches(0.5))
    tf = header.text_frame
    tf.text = text
    paragraph = tf.paragraphs[0]
    paragraph.font.size = Pt(16)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(195, 201, 213)
    return header


def add_footer(slide, number: int | None, label: str) -> List[Any]:
    footer = slide.shapes.add_textbox(CONTENT_LEFT, SLIDE_HEIGHT - Inches(0.7), CONTENT_WIDTH - Inches(1.0), Inches(0.4))
    tf = footer.text_frame
    tf.text = label
    paragraph = tf.paragraphs[0]
    paragraph.font.size = Pt(14)
    paragraph.font.color.rgb = MUTED
    paragraph.font.bold = True

    page_box = slide.shapes.add_textbox(SLIDE_WIDTH - Inches(1.3), SLIDE_HEIGHT - Inches(0.7), Inches(0.9), Inches(0.4))
    page_tf = page_box.text_frame
    page_tf.text = f"{number:02d}" if number else ""
    para = page_tf.paragraphs[0]
    para.font.size = Pt(16)
    para.font.color.rgb = MUTED

    return [footer, page_box]


def add_badge(slide, number: int, color: RGBColor):
    badge = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.35), Inches(1.05), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.text_frame.text = f"{number:02d}"
    p = badge.text_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY if color != SECONDARY else RGBColor(60, 40, 20)
    return badge


def populate_bullets(text_box, bullets: Sequence[str], font_size: int = 20) -> None:
    tf = text_box.text_frame
    tf.clear()
    for index, bullet in enumerate(bullets):
        paragraph = tf.add_paragraph() if index else tf.paragraphs[0]
        paragraph.text = bullet
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = TEXT
        paragraph.level = 0
        paragraph.space_after = Pt(4)


def add_title_slide(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    animated: List[Any] = []
    animated.extend(style_background(slide, accent="right"))
    header = add_header(slide, "Eco-Matic Midterm Showcase")
    animated.append(header)

    title_box = slide.shapes.add_textbox(CONTENT_LEFT, Inches(1.3), CONTENT_WIDTH, Inches(2.5))
    tf = title_box.text_frame
    tf.text = "Eco-Matic Vending Machine"
    title_para = tf.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY

    subtitle = tf.add_paragraph()
    subtitle.text = "Console Application in C#"
    subtitle.font.size = Pt(24)
    subtitle.font.color.rgb = TEXT
    subtitle.space_before = Pt(6)

    footer_line = tf.add_paragraph()
    footer_line.text = "Seanix Real · Sustainable Tech Midterm"
    footer_line.font.size = Pt(18)
    footer_line.font.color.rgb = MUTED
    animated.append(title_box)

    accent_bar = slide.shapes.add_shape(SHAPE.RECTANGLE, CONTENT_LEFT, Inches(4.2), Inches(2.2), Inches(0.12))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    animated.append(accent_bar)

    blurb = slide.shapes.add_textbox(CONTENT_LEFT, Inches(4.5), CONTENT_WIDTH, Inches(1.2))
    blurb_tf = blurb.text_frame
    blurb_tf.text = "Simple, data-backed walkthrough of our SDG-focused vending concept."
    blurb_tf.paragraphs[0].font.size = Pt(20)
    blurb_tf.paragraphs[0].font.color.rgb = TEXT
    animated.append(blurb)

    animated.extend(add_footer(slide, None, "Ready for live demo + documentation"))
    apply_transition(slide)
    apply_object_animations(slide, animated)


def add_section_slide(prs, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = "left" if data["number"] % 2 == 0 else "right"
    animated: List[Any] = []
    animated.extend(style_background(slide, accent=accent))
    header = add_header(slide)
    animated.append(header)
    badge = add_badge(slide, data["number"], ACCENT if accent == "left" else SECONDARY)
    animated.append(badge)

    title_box = slide.shapes.add_textbox(CONTENT_LEFT, Inches(0.95), CONTENT_WIDTH, Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.text = data["title"]
    title_para = title_tf.paragraphs[0]
    title_para.font.size = Pt(34)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY

    subtitle_para = title_tf.add_paragraph()
    subtitle_para.text = data["subtitle"]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = MUTED
    animated.append(title_box)

    bullets_box = slide.shapes.add_textbox(CONTENT_LEFT, Inches(2.1), CONTENT_WIDTH, SLIDE_HEIGHT - Inches(3.0))
    populate_bullets(bullets_box, data["bullets"], font_size=21)
    animated.append(bullets_box)

    animated.extend(add_footer(slide, data["number"], data["title"]))
    apply_transition(slide)
    apply_object_animations(slide, animated)


def add_dual_slide(prs, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    animated: List[Any] = []
    animated.extend(style_background(slide, accent="right"))
    header = add_header(slide)
    animated.append(header)
    badge = add_badge(slide, data["number"], SECONDARY)
    animated.append(badge)

    title_box = slide.shapes.add_textbox(CONTENT_LEFT, Inches(0.95), CONTENT_WIDTH, Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.text = data["title"]
    title_para = title_tf.paragraphs[0]
    title_para.font.size = Pt(34)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY

    subtitle_para = title_tf.add_paragraph()
    subtitle_para.text = data["subtitle"]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = MUTED
    animated.append(title_box)

    card_top = Inches(2.0)
    card_height = SLIDE_HEIGHT - card_top - Inches(1.2)
    column_width = (CONTENT_WIDTH - Inches(0.8)) / 2

    left_box = slide.shapes.add_textbox(CONTENT_LEFT, card_top, column_width, card_height)
    left_tf = left_box.text_frame
    left_tf.text = data["left_header"]
    left_header_para = left_tf.paragraphs[0]
    left_header_para.font.size = Pt(24)
    left_header_para.font.bold = True
    left_header_para.font.color.rgb = ACCENT
    for item in data["left_items"]:
        para = left_tf.add_paragraph()
        para.text = item
        para.font.size = Pt(19)
        para.font.color.rgb = TEXT
        para.space_after = Pt(4)

    right_box = slide.shapes.add_textbox(CONTENT_LEFT + column_width + Inches(0.4), card_top, column_width, card_height)
    right_tf = right_box.text_frame
    right_tf.text = data["right_header"]
    right_header_para = right_tf.paragraphs[0]
    right_header_para.font.size = Pt(24)
    right_header_para.font.bold = True
    right_header_para.font.color.rgb = SECONDARY
    for item in data["right_items"]:
        para = right_tf.add_paragraph()
        para.text = item
        para.font.size = Pt(19)
        para.font.color.rgb = TEXT
        para.space_after = Pt(4)

    animated.extend([left_box, right_box])
    animated.extend(add_footer(slide, data["number"], data["title"]))
    apply_transition(slide)
    apply_object_animations(slide, animated)


def build_presentation(output: Path = OUTPUT) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs)
    for section in SECTION_CONTENT:
        add_section_slide(prs, section)
    add_dual_slide(prs, APPENDIX_CONTENT)

    prs.save(output)


if __name__ == "__main__":
    build_presentation()
